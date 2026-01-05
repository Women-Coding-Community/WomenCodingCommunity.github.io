from collections import Counter
from pathlib import Path
import platform
import hashlib
from datetime import datetime

from pptx import Presentation
import os
import json
import sys
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
import qrcode
from io import BytesIO

# Check platform and conditionally import comtypes (Windows-only for PDF conversion)
IS_WINDOWS = platform.system() == 'Windows'
POWERPOINT_AVAILABLE = False

if IS_WINDOWS:
    try:
        import comtypes.client
        POWERPOINT_AVAILABLE = True
    except ImportError:
        print("Warning: comtypes not available. PDF conversion will be skipped.")
else:
    print(f"Warning: PDF conversion is only supported on Windows. Current platform: {platform.system()}")

# Initialize global powerpoint variable
powerpoint = None

def load_config(config_path="config.json"):
    with open(config_path, 'r', encoding='utf-8') as f:
        return json.load(f)

def generate_certificate_id(name, cert_type, issue_date):
    """Generate a unique certificate ID based on name, type, and date."""
    data = f"{name}|{cert_type}|{issue_date}"
    hash_obj = hashlib.sha256(data.encode('utf-8'))
    return hash_obj.hexdigest()[:12].upper()

def load_certificate_registry(registry_path="../data/output/certificate_registry.json"):
    """Load existing certificate registry or create new one."""
    if os.path.exists(registry_path):
        with open(registry_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    return {"certificates": []}

def save_certificate_registry(registry, registry_path="../data/output/certificate_registry.json"):
    """Save certificate registry to file."""
    os.makedirs(os.path.dirname(registry_path), exist_ok=True)
    with open(registry_path, 'w', encoding='utf-8') as f:
        json.dump(registry, f, indent=2, ensure_ascii=False)

def add_to_registry(registry, cert_id, name, cert_type, issue_date):
    """Add a certificate record to the registry."""
    certificate = {
        "id": cert_id,
        "name": name,
        "type": cert_type,
        "issue_date": issue_date,
        "verification_url": f"https://www.womencodingcommunity.com/verify?cert={cert_id}"
    }

    # Check if certificate already exists
    existing = next((c for c in registry["certificates"] if c["id"] == cert_id), None)
    if not existing:
        registry["certificates"].append(certificate)

    return certificate

def generate_qr_code(verification_url):
    """Generate QR code image for verification URL."""
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=2,
    )
    qr.add_data(verification_url)
    qr.make(fit=True)

    img = qr.make_image(fill_color="black", back_color="white")

    # Save to BytesIO
    img_buffer = BytesIO()
    img.save(img_buffer, format='PNG')
    img_buffer.seek(0)

    return img_buffer

def check_duplicates(names, cert_type):
    counts = Counter(names)
    duplicates = [name for name, count in counts.items() if count > 1]

    if duplicates:
        print(f"\nWARNING: Found {len(duplicates)} duplicate name(s) in {cert_type} list:")
        for name in duplicates:
            print(f"  - {name} (appears {counts[name]} times)")

def load_names(names_file, cert_type):
    all_names = []
    with open(names_file, 'r', encoding='utf-8') as f:
        all_names += [line.strip() for line in f if line.strip()]
    check_duplicates(all_names, cert_type)
    return set(all_names)

def generate_certificates_for_type(names, cert_config, file_type, registry=None, issue_date=None):
    template = cert_config['template']
    input_dir = cert_config["ppt_dir"] if file_type == "pdf" else None
    output_dir = cert_config["ppt_dir"] if file_type == "pptx" else cert_config[
        'pdf_dir']
    placeholder_text = cert_config['placeholder_text']
    cert_type = cert_config['type']

    # Get QR code position parameters (optional)
    qr_left_cm = cert_config.get('qr_left_cm')
    qr_top_cm = cert_config.get('qr_top_cm')
    qr_width_cm = cert_config.get('qr_width_cm', 3.0)
    qr_height_cm = cert_config.get('qr_height_cm', 3.0)

    os.makedirs(output_dir, exist_ok=True)

    print(f"Generating {cert_type.upper()} {cert_type} certificates at {output_dir}")

    file_count = 0

    for i, name in enumerate(names, 1):
        file_name = None
        try:
            if file_type == "pptx":
                file_name = generate_pptx(name, output_dir, placeholder_text,
                                      template, cert_type, registry, issue_date,
                                      qr_left_cm, qr_top_cm, qr_width_cm, qr_height_cm)
            elif file_type == "pdf":
                file_name = generate_pdf(name, input_dir, output_dir)

            print(f"[{i}/{len(names)}] Generated: {file_name}")
            file_count += 1

        except Exception as e:
            print(f"[{i}/{len(names)}] ERROR generating {file_type} "
                  f"certificate for {name}: {e}")

    print(f"\nSuccessfully generated {file_count}/{len(names)} {cert_type} certificates")
    return file_count


def generate_pptx(name, output_dir, placeholder_text, template, cert_type=None,
                  registry=None, issue_date=None, qr_left_cm=None, qr_top_cm=None,
                  qr_width_cm=3.0, qr_height_cm=3.0):
    try:
        prs = Presentation(template)

        # Generate certificate ID and add to registry if provided
        cert_id = None
        verification_url = None
        if registry is not None and issue_date is not None and cert_type is not None:
            cert_id = generate_certificate_id(name, cert_type, issue_date)
            cert_record = add_to_registry(registry, cert_id, name, cert_type, issue_date)
            verification_url = cert_record['verification_url']

        for slide in prs.slides:
            for shape in slide.shapes:
                # Replace name placeholder
                if shape.has_text_frame and shape.text.strip() == placeholder_text:
                    tf = shape.text_frame

                    # Preserve all original formatting from the placeholder
                    original_font = None
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        original_run = tf.paragraphs[0].runs[0]
                        original_font = {
                            'name': original_run.font.name,
                            'size': original_run.font.size,
                            'bold': original_run.font.bold,
                            'italic': original_run.font.italic,
                            'underline': original_run.font.underline,
                            'color': original_run.font.color.rgb if original_run.font.color.type is not None else None
                        }

                    tf.clear()

                    p = tf.paragraphs[0]
                    run = p.add_run()
                    run.text = name

                    # Apply all preserved formatting
                    if original_font:
                        if original_font['name']:
                            run.font.name = original_font['name']
                        if original_font['size']:
                            run.font.size = original_font['size']
                        if original_font['bold'] is not None:
                            run.font.bold = original_font['bold']
                        if original_font['italic'] is not None:
                            run.font.italic = original_font['italic']
                        if original_font['underline'] is not None:
                            run.font.underline = original_font['underline']
                        if original_font['color'] is not None:
                            run.font.color.rgb = original_font['color']

            # Add QR code to slide if verification URL exists
            if verification_url:
                qr_img = generate_qr_code(verification_url)

                # Use configured position (in cm) or default to top right corner
                if qr_left_cm is not None and qr_top_cm is not None:
                    left = Cm(qr_left_cm)
                    top = Cm(qr_top_cm)
                    width = Cm(qr_width_cm)
                    height = Cm(qr_height_cm)
                else:
                    # Default position (top right corner)
                    left = prs.slide_width - Inches(1.5)
                    top = Inches(0.3)
                    width = Inches(1.2)
                    height = Inches(1.2)

                slide.shapes.add_picture(qr_img, left, top, width, height)

        pptx_path = os.path.join(output_dir, f"{name}.pptx")
        prs.save(pptx_path)
        return pptx_path
    except Exception as e:
        raise e

def generate_pdf(name, input_dir, output_path):
    try:
        if not POWERPOINT_AVAILABLE:
            raise RuntimeError("PDF conversion is not available on this platform. Only Windows with PowerPoint is supported.")

        output_path = Path(output_path)

        pptx_path = os.path.abspath(os.path.join(input_dir, f"{name}.pptx"))

        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PPTX not found: {pptx_path}")

        presentation = powerpoint.Presentations.Open(pptx_path)

        pdf_path = output_path / f"{name}.pdf"

        presentation.SaveAs(str(pdf_path.resolve()), 32)

        presentation.Close()

        return pdf_path

    except Exception as e:
        raise e

def check_metrics(names, cert_config, file_type):
    cert_type = cert_config['type']
    output_dir = cert_config["ppt_dir"] if file_type == "pptx" else \
        cert_config['pdf_dir']

    print(f"Checking metrics {cert_type.upper()} certificates")

    folder_path = Path(output_dir)

    if not folder_path.exists():
        print(f"ERROR: Directory does not exist: {output_dir}")
        return 0, len(names)

    existing_files = {f.stem for f in folder_path.glob(f"*.{file_type}")}

    print(f"\nExpected certificates: {len(names)}")
    print(f"Found certificates: {len(existing_files)}")

    missing = []
    for name in names:
        if name not in existing_files:
            missing.append(name)

    if missing:
        print(f"\nMissing {len(missing)} certificate(s):")
        for name in missing:
            print(f"  - {name}")
    else:
        print(f"\nAll {cert_type} certificates are present!")

def main():
    global powerpoint
    try:
        config = load_config()

        # Load certificate registry
        registry = load_certificate_registry()

        # Get current issue date
        issue_date = datetime.now().strftime("%Y-%m-%d")

        # Initialize PowerPoint if available
        if POWERPOINT_AVAILABLE:
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 1
        else:
            print("\nPDF generation will be skipped (not available on this platform).\n")

        for cert_config in config['certificate_types']:
            names_file = cert_config['names_file']
            cert_type = cert_config['type']

            names = load_names(names_file, cert_type)

            # Generate PPTX certificates
            pptx_generated = generate_certificates_for_type(names,
                                                            cert_config,
                                                            "pptx",
                                                            registry,
                                                            issue_date)
            check_metrics(names, cert_config, "pptx")

            # Generate PDF certificates only if PowerPoint is available
            pdf_generated = 0
            if POWERPOINT_AVAILABLE:
                pdf_generated = generate_certificates_for_type(names,
                                                               cert_config, "pdf")
                check_metrics(names, cert_config, "pdf")
            else:
                print(f"Skipping PDF generation for {cert_type} (PowerPoint not available)")

            total_certificates = len(names)
            print(f"Type: {cert_config['type']} Total: {total_certificates} "
                  f"PPTX Generated: {pptx_generated} PDF Generated: {pdf_generated}")

        # Save updated registry
        save_certificate_registry(registry)
        print(f"\nCertificate registry saved with {len(registry['certificates'])} total certificates")


    except Exception as e:
        print(f"Error while running the certificate generation automation:"
              f" {e}")
        return 1
    finally:
        if powerpoint:
            powerpoint.Quit()

if __name__ == "__main__":
    sys.exit(main())
