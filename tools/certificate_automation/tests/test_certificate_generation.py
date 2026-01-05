import unittest
import sys
import os
import tempfile
import shutil
from unittest.mock import Mock, patch, MagicMock
from pptx import Presentation
from pptx.util import Inches, Pt

# Add src directory to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from generate_certificates import (
    generate_pptx,
    check_duplicates,
    load_names
)


class TestCertificateGeneration(unittest.TestCase):
    """Test cases for certificate generation functionality."""

    def setUp(self):
        """Set up test fixtures."""
        self.test_dir = tempfile.mkdtemp()
        self.output_dir = os.path.join(self.test_dir, 'output')
        os.makedirs(self.output_dir, exist_ok=True)

        # Create a simple test template
        self.template_path = os.path.join(self.test_dir, 'template.pptx')
        self._create_test_template()

    def tearDown(self):
        """Tear down test fixtures."""
        shutil.rmtree(self.test_dir)

    def _create_test_template(self):
        """Create a simple PPTX template for testing."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add a text box with placeholder text
        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Sample Sample"

        prs.save(self.template_path)

    def test_generate_pptx_creates_file(self):
        """Test that generate_pptx creates output file."""
        name = "John Smith"
        registry = {"certificates": []}
        issue_date = "2026-01-04"

        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="mentee",
            registry=registry,
            issue_date=issue_date
        )

        # Verify file was created
        self.assertTrue(os.path.exists(result))
        self.assertTrue(result.endswith('.pptx'))

    def test_generate_pptx_replaces_placeholder(self):
        """Test that placeholder text is replaced with name."""
        name = "Jane Doe"
        registry = {"certificates": []}
        issue_date = "2026-01-04"

        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="mentor",
            registry=registry,
            issue_date=issue_date
        )

        # Open the generated file and check content
        prs = Presentation(result)
        slide = prs.slides[0]

        # Find text in shapes
        found_name = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                if name in shape.text:
                    found_name = True
                    break

        self.assertTrue(found_name, f"Name '{name}' not found in generated certificate")

    def test_generate_pptx_adds_to_registry(self):
        """Test that certificate is added to registry."""
        name = "Alice Johnson"
        registry = {"certificates": []}
        issue_date = "2026-01-04"

        generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="volunteer",
            registry=registry,
            issue_date=issue_date
        )

        # Verify certificate was added to registry
        self.assertEqual(len(registry['certificates']), 1)
        self.assertEqual(registry['certificates'][0]['name'], name)
        self.assertEqual(registry['certificates'][0]['type'], 'volunteer')

    def test_generate_pptx_adds_qr_code(self):
        """Test that QR code is added to certificate."""
        name = "Bob Smith"
        registry = {"certificates": []}
        issue_date = "2026-01-04"

        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="mentee",
            registry=registry,
            issue_date=issue_date
        )

        # Open the generated file
        prs = Presentation(result)
        slide = prs.slides[0]

        # Count image shapes (QR code should be one)
        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)  # 13 = PICTURE

        self.assertGreater(image_count, 0, "No QR code image found in certificate")

    def test_generate_pptx_without_registry(self):
        """Test that generate_pptx works without registry (backwards compatibility)."""
        name = "Legacy User"

        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path
        )

        # Should create file without error
        self.assertTrue(os.path.exists(result))

    def test_check_duplicates_no_duplicates(self):
        """Test check_duplicates with unique names."""
        names = ["John Smith", "Jane Doe", "Alice Johnson"]

        # Should run without raising exception
        try:
            check_duplicates(names, "mentee")
        except Exception as e:
            self.fail(f"check_duplicates raised exception: {e}")

    def test_check_duplicates_with_duplicates(self):
        """Test check_duplicates identifies duplicates."""
        names = ["John Smith", "Jane Doe", "John Smith", "Alice Johnson"]

        # Capture output
        import io
        from contextlib import redirect_stdout

        f = io.StringIO()
        with redirect_stdout(f):
            check_duplicates(names, "mentee")

        output = f.getvalue()

        # Should print warning about duplicates
        self.assertIn("WARNING", output)
        self.assertIn("John Smith", output)

    def test_load_names_from_file(self):
        """Test loading names from file."""
        # Create test names file
        names_file = os.path.join(self.test_dir, 'names.txt')
        with open(names_file, 'w', encoding='utf-8') as f:
            f.write("John Smith\n")
            f.write("Jane Doe\n")
            f.write("Alice Johnson\n")

        names = load_names(names_file, "mentee")

        self.assertEqual(len(names), 3)
        self.assertIn("John Smith", names)
        self.assertIn("Jane Doe", names)
        self.assertIn("Alice Johnson", names)

    def test_load_names_removes_duplicates(self):
        """Test that load_names returns unique names."""
        # Create test names file with duplicates
        names_file = os.path.join(self.test_dir, 'names_dup.txt')
        with open(names_file, 'w', encoding='utf-8') as f:
            f.write("John Smith\n")
            f.write("Jane Doe\n")
            f.write("John Smith\n")  # Duplicate

        names = load_names(names_file, "mentee")

        # Should return set with unique names
        self.assertEqual(len(names), 2)

    def test_load_names_strips_whitespace(self):
        """Test that load_names strips whitespace from names."""
        # Create test names file with extra whitespace
        names_file = os.path.join(self.test_dir, 'names_ws.txt')
        with open(names_file, 'w', encoding='utf-8') as f:
            f.write("  John Smith  \n")
            f.write("\tJane Doe\n")
            f.write("Alice Johnson   \n")

        names = load_names(names_file, "mentee")

        self.assertIn("John Smith", names)
        self.assertIn("Jane Doe", names)
        self.assertIn("Alice Johnson", names)

    def test_load_names_skips_empty_lines(self):
        """Test that load_names skips empty lines."""
        # Create test names file with empty lines
        names_file = os.path.join(self.test_dir, 'names_empty.txt')
        with open(names_file, 'w', encoding='utf-8') as f:
            f.write("John Smith\n")
            f.write("\n")
            f.write("Jane Doe\n")
            f.write("   \n")
            f.write("Alice Johnson\n")

        names = load_names(names_file, "mentee")

        self.assertEqual(len(names), 3)


if __name__ == '__main__':
    unittest.main()
