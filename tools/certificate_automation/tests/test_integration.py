import unittest
import sys
import os
import tempfile
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt

# Add src directory to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))

from generate_certificates import (
    generate_certificate_id,
    generate_qr_code,
    load_certificate_registry,
    save_certificate_registry,
    add_to_registry,
    generate_pptx
)


class TestIntegration(unittest.TestCase):
    """Integration tests for complete certificate generation workflow."""

    def setUp(self):
        """Set up test fixtures."""
        self.test_dir = tempfile.mkdtemp()
        self.output_dir = os.path.join(self.test_dir, 'output')
        os.makedirs(self.output_dir, exist_ok=True)

        self.template_path = os.path.join(self.test_dir, 'template.pptx')
        self.registry_path = os.path.join(self.test_dir, 'registry.json')

        self._create_test_template()

    def tearDown(self):
        """Tear down test fixtures."""
        shutil.rmtree(self.test_dir)

    def _create_test_template(self):
        """Create a simple PPTX template for testing."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(2)
        top = Inches(3)
        width = Inches(6)
        height = Inches(1)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "Sample Sample"

        prs.save(self.template_path)

    def test_complete_certificate_workflow(self):
        """Test complete workflow: generate ID, create QR, add to registry, create PPTX."""
        name = "John Smith"
        cert_type = "mentee"
        issue_date = "2026-01-04"

        # Step 1: Load/create registry
        registry = load_certificate_registry(self.registry_path)
        self.assertEqual(len(registry['certificates']), 0)

        # Step 2: Generate certificate
        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type=cert_type,
            registry=registry,
            issue_date=issue_date
        )

        # Step 3: Verify PPTX was created
        self.assertTrue(os.path.exists(result))

        # Step 4: Verify registry was updated
        self.assertEqual(len(registry['certificates']), 1)
        cert_record = registry['certificates'][0]

        self.assertEqual(cert_record['name'], name)
        self.assertEqual(cert_record['type'], cert_type)
        self.assertEqual(cert_record['issue_date'], issue_date)
        self.assertIn('id', cert_record)
        self.assertIn('verification_url', cert_record)

        # Step 5: Save registry
        save_certificate_registry(registry, self.registry_path)

        # Step 6: Verify registry file was created
        self.assertTrue(os.path.exists(self.registry_path))

        # Step 7: Load registry again and verify
        loaded_registry = load_certificate_registry(self.registry_path)
        self.assertEqual(len(loaded_registry['certificates']), 1)
        self.assertEqual(loaded_registry['certificates'][0]['name'], name)

    def test_multiple_certificates_workflow(self):
        """Test generating multiple certificates in one batch."""
        names = ["Alice Johnson", "Bob Smith", "Carol White"]
        cert_type = "mentor"
        issue_date = "2026-01-04"

        registry = load_certificate_registry(self.registry_path)

        # Generate certificates for all names
        for name in names:
            generate_pptx(
                name=name,
                output_dir=self.output_dir,
                placeholder_text="Sample Sample",
                template=self.template_path,
                cert_type=cert_type,
                registry=registry,
                issue_date=issue_date
            )

        # Verify all certificates were created
        generated_files = os.listdir(self.output_dir)
        self.assertEqual(len(generated_files), 3)

        # Verify all certificates are in registry
        self.assertEqual(len(registry['certificates']), 3)

        # Verify each name is in registry
        registry_names = [cert['name'] for cert in registry['certificates']]
        for name in names:
            self.assertIn(name, registry_names)

        # Verify all IDs are unique
        cert_ids = [cert['id'] for cert in registry['certificates']]
        self.assertEqual(len(cert_ids), len(set(cert_ids)))

    def test_different_certificate_types(self):
        """Test generating different types of certificates."""
        test_cases = [
            ("User1", "mentee", "2026-01-04"),
            ("User2", "mentor", "2026-01-04"),
            ("User3", "volunteer", "2026-01-04"),
            ("User4", "leader", "2026-01-04")
        ]

        registry = load_certificate_registry(self.registry_path)

        for name, cert_type, issue_date in test_cases:
            generate_pptx(
                name=name,
                output_dir=self.output_dir,
                placeholder_text="Sample Sample",
                template=self.template_path,
                cert_type=cert_type,
                registry=registry,
                issue_date=issue_date
            )

        # Verify all types are present in registry
        cert_types = [cert['type'] for cert in registry['certificates']]
        self.assertIn("mentee", cert_types)
        self.assertIn("mentor", cert_types)
        self.assertIn("volunteer", cert_types)
        self.assertIn("leader", cert_types)

    def test_same_name_different_types_different_ids(self):
        """Test that same name with different types gets different IDs."""
        name = "John Smith"
        issue_date = "2026-01-04"

        registry = load_certificate_registry(self.registry_path)

        # Generate mentee certificate
        generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="mentee",
            registry=registry,
            issue_date=issue_date
        )

        # Generate mentor certificate for same person
        generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type="mentor",
            registry=registry,
            issue_date=issue_date
        )

        # Should have 2 certificates with different IDs
        self.assertEqual(len(registry['certificates']), 2)

        id1 = registry['certificates'][0]['id']
        id2 = registry['certificates'][1]['id']

        self.assertNotEqual(id1, id2)

    def test_qr_code_in_generated_certificate(self):
        """Test that QR code is actually embedded in the generated certificate."""
        name = "QR Test User"
        cert_type = "mentee"
        issue_date = "2026-01-04"

        registry = load_certificate_registry(self.registry_path)

        result = generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type=cert_type,
            registry=registry,
            issue_date=issue_date
        )

        # Open generated certificate
        prs = Presentation(result)
        slide = prs.slides[0]

        # Count images (should have QR code)
        image_count = sum(1 for shape in slide.shapes if shape.shape_type == 13)

        self.assertGreater(image_count, 0, "No QR code found in certificate")

        # Verify certificate is in registry
        cert_id = registry['certificates'][0]['id']
        self.assertIsNotNone(cert_id)
        self.assertEqual(len(cert_id), 12)

    def test_verification_url_contains_cert_id(self):
        """Test that verification URL contains the correct certificate ID."""
        name = "URL Test User"
        cert_type = "mentor"
        issue_date = "2026-01-04"

        registry = load_certificate_registry(self.registry_path)

        generate_pptx(
            name=name,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type=cert_type,
            registry=registry,
            issue_date=issue_date
        )

        cert = registry['certificates'][0]
        cert_id = cert['id']
        verification_url = cert['verification_url']

        # Verify URL contains certificate ID
        self.assertIn(cert_id, verification_url)
        self.assertIn("womencodingcommunity.com", verification_url)
        self.assertIn("verify?cert=", verification_url)

    def test_registry_persistence(self):
        """Test that registry persists across multiple save/load cycles."""
        name1 = "Persistent User 1"
        name2 = "Persistent User 2"
        cert_type = "mentee"
        issue_date = "2026-01-04"

        # First batch
        registry = load_certificate_registry(self.registry_path)
        generate_pptx(
            name=name1,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type=cert_type,
            registry=registry,
            issue_date=issue_date
        )
        save_certificate_registry(registry, self.registry_path)

        # Second batch - load existing registry
        registry = load_certificate_registry(self.registry_path)
        self.assertEqual(len(registry['certificates']), 1)

        generate_pptx(
            name=name2,
            output_dir=self.output_dir,
            placeholder_text="Sample Sample",
            template=self.template_path,
            cert_type=cert_type,
            registry=registry,
            issue_date=issue_date
        )
        save_certificate_registry(registry, self.registry_path)

        # Load final registry
        final_registry = load_certificate_registry(self.registry_path)
        self.assertEqual(len(final_registry['certificates']), 2)

        names = [cert['name'] for cert in final_registry['certificates']]
        self.assertIn(name1, names)
        self.assertIn(name2, names)


if __name__ == '__main__':
    unittest.main()
